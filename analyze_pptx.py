#!/usr/bin/env python3
"""Analyze and compare PPTX files to understand their design."""

import io
import sys
from collections import Counter
from pathlib import Path
from zipfile import ZipFile

from lxml import etree
from pptx import Presentation as PptxPresentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

NAMESPACES = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def analyze_pptx(filepath: str, label: str = ""):
    """Analyze a PPTX file and print its design characteristics."""
    print(f"\n{'='*60}")
    print(f"ANALYZING: {label or filepath}")
    print(f"{'='*60}")

    prs = PptxPresentation(filepath)

    # Basic info
    print(f"\nSlide dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    print(f"Number of slides: {len(prs.slides)}")
    print(f"Number of layouts: {len(prs.slide_layouts)}")

    # List layouts
    print("\nSlide Layouts:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  {i}: {layout.name}")

    # Extract theme colors
    print("\nTheme Colors:")
    try:
        slide_master = prs.slide_master
        slide_master_part = slide_master.part
        theme_part = slide_master_part.part_related_by(RT.THEME)
        theme = etree.fromstring(theme_part.blob)

        color_elements = ["dk1", "lt1", "dk2", "lt2",
                         "accent1", "accent2", "accent3", "accent4", "accent5", "accent6"]
        for color_name in color_elements:
            xpath = f".//a:clrScheme/a:{color_name}"
            elements = theme.xpath(xpath, namespaces=NAMESPACES)
            if elements:
                srgb = elements[0].xpath("a:srgbClr/@val", namespaces=NAMESPACES)
                if srgb:
                    print(f"  {color_name}: #{srgb[0]}")
                else:
                    sys_clr = elements[0].xpath("a:sysClr/@lastClr", namespaces=NAMESPACES)
                    if sys_clr:
                        print(f"  {color_name}: #{sys_clr[0]} (system)")
    except Exception as e:
        print(f"  Error extracting theme colors: {e}")

    # Extract fonts from theme
    print("\nTheme Fonts:")
    try:
        major_latin = theme.xpath(
            ".//a:fontScheme/a:majorFont/a:latin/@typeface",
            namespaces=NAMESPACES,
        )
        if major_latin:
            print(f"  Major (heading): {major_latin[0]}")

        minor_latin = theme.xpath(
            ".//a:fontScheme/a:minorFont/a:latin/@typeface",
            namespaces=NAMESPACES,
        )
        if minor_latin:
            print(f"  Minor (body): {minor_latin[0]}")
    except Exception as e:
        print(f"  Error extracting fonts: {e}")

    # Extract fonts actually used in shapes
    print("\nFonts Used in Shapes:")
    fonts = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.font.name:
                            fonts[para.font.name] += 1
                        for run in para.runs:
                            if run.font.name:
                                fonts[run.font.name] += 1
            except:
                pass
    for font, count in fonts.most_common(10):
        print(f"  {font}: {count} uses")

    # Analyze each slide
    print("\nSlide Analysis:")
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        print(f"\n  Slide {i+1} (Layout: {layout_name}):")

        # Count shapes
        shape_types = Counter()
        for shape in slide.shapes:
            shape_type = str(type(shape).__name__)
            shape_types[shape_type] += 1

        print(f"    Shapes: {dict(shape_types)}")

        # Check background
        try:
            bg = slide.background
            if bg.fill.type is not None:
                fill_type = str(bg.fill.type).replace("MSO_FILL_TYPE.", "")
                print(f"    Background: {fill_type}")
                if 'SOLID' in str(bg.fill.type):
                    try:
                        rgb = bg.fill.fore_color.rgb
                        if rgb:
                            print(f"    Background Color: #{rgb}")
                    except:
                        pass
        except:
            pass

        # Analyze text content and styling
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) < 100:
                    # Get first paragraph's formatting
                    try:
                        p = shape.text_frame.paragraphs[0]
                        font_info = []
                        if p.font.name:
                            font_info.append(f"font={p.font.name}")
                        if p.font.size:
                            font_info.append(f"size={p.font.size.pt}pt")
                        if p.font.bold:
                            font_info.append("bold")
                        if p.font.color and p.font.color.rgb:
                            font_info.append(f"color=#{p.font.color.rgb}")
                        info = ", ".join(font_info) if font_info else "default"
                        print(f"    Text: '{text[:50]}...' [{info}]")
                    except:
                        print(f"    Text: '{text[:50]}...'")

    # Extract colors from shapes
    print("\nColors Used in Shapes (by frequency):")
    fill_colors = Counter()
    text_colors = Counter()

    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    if 'SOLID' in str(shape.fill.type):
                        if shape.fill.fore_color and shape.fill.fore_color.type is not None:
                            if 'RGB' in str(shape.fill.fore_color.type):
                                fill_colors[f"#{shape.fill.fore_color.rgb}"] += 1
            except:
                pass

            try:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.color and run.font.color.type is not None:
                                if 'RGB' in str(run.font.color.type):
                                    text_colors[f"#{run.font.color.rgb}"] += 1
            except:
                pass

    print("  Fill colors:")
    for color, count in fill_colors.most_common(10):
        print(f"    {color}: {count}")
    print("  Text colors:")
    for color, count in text_colors.most_common(10):
        print(f"    {color}: {count}")


def compare_pptx(template_path: str, generated_path: str):
    """Compare template and generated PPTX files."""
    print("\n" + "="*60)
    print("COMPARISON SUMMARY")
    print("="*60)

    template = PptxPresentation(template_path)
    generated = PptxPresentation(generated_path)

    print(f"\n{'Attribute':<30} {'Template':<20} {'Generated':<20}")
    print("-"*70)
    print(f"{'Slides':<30} {len(template.slides):<20} {len(generated.slides):<20}")
    print(f"{'Layouts':<30} {len(template.slide_layouts):<20} {len(generated.slide_layouts):<20}")
    print(f"{'Width (inches)':<30} {template.slide_width.inches:<20.2f} {generated.slide_width.inches:<20.2f}")
    print(f"{'Height (inches)':<30} {template.slide_height.inches:<20.2f} {generated.slide_height.inches:<20.2f}")

    # Compare layouts used
    print("\nLayouts used in generated presentation:")
    for i, slide in enumerate(generated.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        print(f"  Slide {i+1}: {layout_name}")


if __name__ == "__main__":
    template_path = "/root/fyc/output/58388ab2/uploaded_template_20200417_Scoping Präsentation.pptx"
    generated_path = "/root/fyc/output/58388ab2/presentation.pptx"

    analyze_pptx(template_path, "UPLOADED TEMPLATE")
    analyze_pptx(generated_path, "GENERATED PRESENTATION")
    compare_pptx(template_path, generated_path)
