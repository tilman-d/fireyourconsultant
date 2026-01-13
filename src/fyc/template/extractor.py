"""PPTX template style extractor using python-pptx."""

import io
from collections import Counter
from pathlib import Path
from typing import Optional
from zipfile import ZipFile

from lxml import etree
from pptx import Presentation as PptxPresentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from ..models import (
    TemplateProfile,
    ThemeColors,
    ThemeFonts,
    BackgroundStyle,
    ExtractedLayout,
    PlaceholderInfo,
    ExtractedColorPalette,
)


class PptxTemplateExtractor:
    """Extracts style information from PPTX files."""

    # XML namespaces for XPath queries
    NAMESPACES = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    def __init__(self, file_content: bytes, filename: str, output_dir: str):
        self.file_content = file_content
        self.filename = filename
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.prs = PptxPresentation(io.BytesIO(file_content))

    def extract(self) -> TemplateProfile:
        """Extract complete template profile from PPTX."""
        theme_colors = self._extract_theme_colors()
        theme_fonts = self._extract_theme_fonts()
        master_background = self._extract_master_background()
        layouts = self._extract_layouts()
        extracted_images = self._extract_background_images()
        # NEW: Extract actual colors from shapes
        extracted_palette = self._extract_shape_colors()

        return TemplateProfile(
            source_file=self.filename,
            theme_colors=theme_colors,
            theme_fonts=theme_fonts,
            master_background=master_background,
            layouts=layouts,
            extracted_images=extracted_images,
            extracted_palette=extracted_palette,
            template_bytes=self.file_content,  # Store for use as base
        )

    def _get_theme_xml(self):
        """Get parsed theme XML from the presentation using lxml."""
        try:
            slide_master = self.prs.slide_master
            slide_master_part = slide_master.part
            theme_part = slide_master_part.part_related_by(RT.THEME)
            # Parse with lxml directly to get proper xpath support
            return etree.fromstring(theme_part.blob)
        except Exception:
            return None

    def _extract_theme_colors(self) -> ThemeColors:
        """Extract theme colors from clrScheme in theme XML."""
        theme = self._get_theme_xml()
        if theme is None:
            return ThemeColors()

        colors = {}

        color_elements = [
            "dk1", "lt1", "dk2", "lt2",
            "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
            "hlink", "folHlink",
        ]

        for color_name in color_elements:
            xpath = f".//a:clrScheme/a:{color_name}"
            elements = theme.xpath(xpath, namespaces=self.NAMESPACES)
            if elements:
                # Try srgbClr first, then sysClr
                srgb = elements[0].xpath("a:srgbClr/@val", namespaces=self.NAMESPACES)
                if srgb:
                    colors[color_name] = f"#{srgb[0]}"
                else:
                    sys_clr = elements[0].xpath(
                        "a:sysClr/@lastClr", namespaces=self.NAMESPACES
                    )
                    if sys_clr:
                        colors[color_name] = f"#{sys_clr[0]}"

        return ThemeColors(**colors)

    def _extract_theme_fonts(self) -> ThemeFonts:
        """Extract fonts - prefer actual fonts from shapes over theme fonts."""
        # First, get fonts actually used in shapes (more accurate)
        shape_fonts = self._extract_fonts_from_shapes()

        # Fall back to theme fonts if no shape fonts found
        theme = self._get_theme_xml()
        theme_major = "Calibri"
        theme_minor = "Calibri"

        if theme is not None:
            major_latin = theme.xpath(
                ".//a:fontScheme/a:majorFont/a:latin/@typeface",
                namespaces=self.NAMESPACES,
            )
            if major_latin:
                theme_major = major_latin[0]

            minor_latin = theme.xpath(
                ".//a:fontScheme/a:minorFont/a:latin/@typeface",
                namespaces=self.NAMESPACES,
            )
            if minor_latin:
                theme_minor = minor_latin[0]

        # Use shape fonts if found, otherwise theme fonts
        return ThemeFonts(
            major_latin=shape_fonts.get('heading', theme_major),
            minor_latin=shape_fonts.get('body', theme_minor),
        )

    def _extract_fonts_from_shapes(self) -> dict:
        """Extract actually used fonts from shapes in slides."""
        from collections import Counter
        fonts = Counter()

        # Check slides
        for slide in self.prs.slides:
            for shape in slide.shapes:
                self._collect_fonts_from_shape(shape, fonts)

        # Check master
        try:
            for shape in self.prs.slide_master.shapes:
                self._collect_fonts_from_shape(shape, fonts)
        except:
            pass

        # Check layouts
        try:
            for layout in self.prs.slide_master.slide_layouts:
                for shape in layout.shapes:
                    self._collect_fonts_from_shape(shape, fonts)
        except:
            pass

        if not fonts:
            return {}

        # Get most common font family (ignore variants like Light, Medium)
        font_families = Counter()
        for font, count in fonts.items():
            # Extract base family name
            base = font.split()[0] if font else ""
            if base:
                font_families[base] += count

        most_common = font_families.most_common(1)
        if most_common:
            base_family = most_common[0][0]
            # Find the most common variant for heading (often Medium/Bold) and body (often Light/Regular)
            heading_font = base_family
            body_font = base_family

            for font in fonts:
                if font.startswith(base_family):
                    if 'Medium' in font or 'Bold' in font:
                        heading_font = font
                    elif 'Light' in font or 'Regular' in font:
                        body_font = font

            # If no specific variants, use the most common
            if heading_font == base_family:
                heading_font = fonts.most_common(1)[0][0]
            if body_font == base_family:
                body_font = fonts.most_common(1)[0][0]

            return {'heading': heading_font, 'body': body_font}

        return {}

    def _collect_fonts_from_shape(self, shape, fonts: Counter):
        """Collect font names from a shape."""
        try:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.font.name:
                        fonts[para.font.name] += 1
                    for run in para.runs:
                        if run.font.name:
                            fonts[run.font.name] += 1
        except:
            pass

    def _extract_theme_fonts_old(self) -> ThemeFonts:
        """Extract theme fonts from fontScheme in theme XML (legacy method)."""
        theme = self._get_theme_xml()
        if theme is None:
            return ThemeFonts()

        fonts = {}

        # Major font (headings)
        major_latin = theme.xpath(
            ".//a:fontScheme/a:majorFont/a:latin/@typeface",
            namespaces=self.NAMESPACES,
        )
        if major_latin:
            fonts["major_latin"] = major_latin[0]

        # Minor font (body)
        minor_latin = theme.xpath(
            ".//a:fontScheme/a:minorFont/a:latin/@typeface",
            namespaces=self.NAMESPACES,
        )
        if minor_latin:
            fonts["minor_latin"] = minor_latin[0]

        # East Asian fonts
        major_ea = theme.xpath(
            ".//a:fontScheme/a:majorFont/a:ea/@typeface",
            namespaces=self.NAMESPACES,
        )
        if major_ea:
            fonts["major_ea"] = major_ea[0]

        minor_ea = theme.xpath(
            ".//a:fontScheme/a:minorFont/a:ea/@typeface",
            namespaces=self.NAMESPACES,
        )
        if minor_ea:
            fonts["minor_ea"] = minor_ea[0]

        # Complex script fonts
        major_cs = theme.xpath(
            ".//a:fontScheme/a:majorFont/a:cs/@typeface",
            namespaces=self.NAMESPACES,
        )
        if major_cs:
            fonts["major_cs"] = major_cs[0]

        minor_cs = theme.xpath(
            ".//a:fontScheme/a:minorFont/a:cs/@typeface",
            namespaces=self.NAMESPACES,
        )
        if minor_cs:
            fonts["minor_cs"] = minor_cs[0]

        return ThemeFonts(**fonts)

    def _extract_master_background(self) -> Optional[BackgroundStyle]:
        """Extract background style from slide master."""
        try:
            slide_master = self.prs.slide_master
            bg = slide_master.background

            if bg.fill.type is not None:
                return self._extract_fill_style(bg.fill)
        except Exception:
            pass
        return None

    def _extract_fill_style(self, fill) -> BackgroundStyle:
        """Convert a FillFormat to BackgroundStyle."""
        try:
            # Determine fill type
            fill_type_str = str(fill.type) if fill.type else "solid"
            fill_type = fill_type_str.replace("MSO_FILL_TYPE.", "").lower()
        except Exception:
            fill_type = "solid"

        style = BackgroundStyle(fill_type=fill_type)

        if fill_type == "solid":
            try:
                if fill.fore_color and fill.fore_color.rgb:
                    style.solid_color = f"#{fill.fore_color.rgb}"
            except Exception:
                pass
        elif fill_type == "gradient":
            # Extract gradient stops
            style.gradient_colors = []
            try:
                for stop in fill.gradient_stops:
                    if stop.color.rgb:
                        style.gradient_colors.append(f"#{stop.color.rgb}")
                style.gradient_angle = fill.gradient_angle or 0
            except Exception:
                pass

        return style

    def _extract_layouts(self) -> list[ExtractedLayout]:
        """Extract all slide layouts with placeholder information."""
        layouts = []

        try:
            for idx, layout in enumerate(self.prs.slide_master.slide_layouts):
                placeholders = []

                for placeholder in layout.placeholders:
                    try:
                        ph_type = str(placeholder.placeholder_format.type)
                        ph_type = ph_type.replace("PP_PLACEHOLDER_TYPE.", "")

                        ph_info = PlaceholderInfo(
                            idx=placeholder.placeholder_format.idx,
                            type=ph_type,
                            left=placeholder.left.inches if placeholder.left else 0,
                            top=placeholder.top.inches if placeholder.top else 0,
                            width=placeholder.width.inches if placeholder.width else 0,
                            height=placeholder.height.inches if placeholder.height else 0,
                        )

                        # Try to extract font info from placeholder
                        if placeholder.has_text_frame:
                            for paragraph in placeholder.text_frame.paragraphs:
                                if paragraph.font.name:
                                    ph_info.font_name = paragraph.font.name
                                if paragraph.font.size:
                                    ph_info.font_size = paragraph.font.size.pt
                                if paragraph.font.bold is not None:
                                    ph_info.font_bold = paragraph.font.bold
                                break

                        placeholders.append(ph_info)
                    except Exception:
                        continue

                # Extract layout background if different from master
                bg_style = None
                try:
                    if layout.background.fill.type is not None:
                        bg_style = self._extract_fill_style(layout.background.fill)
                except Exception:
                    pass

                layouts.append(
                    ExtractedLayout(
                        name=layout.name,
                        idx=idx,
                        placeholders=placeholders,
                        background=bg_style,
                    )
                )
        except Exception:
            pass

        return layouts

    def _extract_shape_colors(self) -> ExtractedColorPalette:
        """Extract actual colors from shapes in slides, master, and layouts.

        This captures the real visual colors used in the template, which are
        often set via explicit RGB values rather than theme references.
        """
        fill_colors = Counter()  # Colors used for shape fills
        text_colors = Counter()  # Colors used for text
        bg_colors = Counter()    # Colors used for backgrounds

        def rgb_to_hex(rgb) -> Optional[str]:
            """Convert RGBColor to hex string."""
            if rgb is None:
                return None
            try:
                if isinstance(rgb, RGBColor):
                    return f"#{rgb}"
                return f"#{rgb}"
            except Exception:
                return None

        def extract_colors_from_shape(shape):
            """Extract fill and text colors from a shape."""
            # Get fill color
            try:
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    fill_type = str(shape.fill.type)
                    if 'SOLID' in fill_type:
                        if shape.fill.fore_color and shape.fill.fore_color.type is not None:
                            color_type = str(shape.fill.fore_color.type)
                            if 'RGB' in color_type:
                                hex_color = rgb_to_hex(shape.fill.fore_color.rgb)
                                if hex_color:
                                    fill_colors[hex_color.upper()] += 1
            except Exception:
                pass

            # Get text colors
            try:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.color and run.font.color.type is not None:
                                color_type = str(run.font.color.type)
                                if 'RGB' in color_type:
                                    hex_color = rgb_to_hex(run.font.color.rgb)
                                    if hex_color:
                                        text_colors[hex_color.upper()] += 1
            except Exception:
                pass

            # Get line/border color
            try:
                if hasattr(shape, 'line') and shape.line.fill.type is not None:
                    line_type = str(shape.line.fill.type)
                    if 'SOLID' in line_type:
                        if shape.line.color and shape.line.color.type is not None:
                            color_type = str(shape.line.color.type)
                            if 'RGB' in color_type:
                                hex_color = rgb_to_hex(shape.line.color.rgb)
                                if hex_color:
                                    fill_colors[hex_color.upper()] += 1
            except Exception:
                pass

        def extract_background_color(bg):
            """Extract background fill color."""
            try:
                if bg.fill.type is not None:
                    fill_type = str(bg.fill.type)
                    if 'SOLID' in fill_type:
                        if bg.fill.fore_color and bg.fill.fore_color.type is not None:
                            color_type = str(bg.fill.fore_color.type)
                            if 'RGB' in color_type:
                                hex_color = rgb_to_hex(bg.fill.fore_color.rgb)
                                if hex_color:
                                    bg_colors[hex_color.upper()] += 1
            except Exception:
                pass

        # Extract from slide master
        try:
            master = self.prs.slide_master
            extract_background_color(master.background)
            for shape in master.shapes:
                extract_colors_from_shape(shape)
        except Exception:
            pass

        # Extract from all layouts
        try:
            for layout in self.prs.slide_master.slide_layouts:
                extract_background_color(layout.background)
                for shape in layout.shapes:
                    extract_colors_from_shape(shape)
        except Exception:
            pass

        # Extract from all slides
        try:
            for slide in self.prs.slides:
                extract_background_color(slide.background)
                for shape in slide.shapes:
                    extract_colors_from_shape(shape)
        except Exception:
            pass

        # Build the palette from extracted colors
        def is_neutral(color: str) -> bool:
            """Check if color is white, black, or near-neutral gray."""
            color = color.upper().lstrip('#')
            if len(color) != 6:
                return True
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            # Pure white/black
            if (r, g, b) in [(255, 255, 255), (0, 0, 0)]:
                return True
            # Near-neutral grays (low saturation)
            max_rgb = max(r, g, b)
            min_rgb = min(r, g, b)
            if max_rgb - min_rgb < 20 and (max_rgb > 200 or min_rgb < 55):
                return True
            return False

        def is_dark(color: str) -> bool:
            """Check if color is dark (good for text)."""
            color = color.upper().lstrip('#')
            if len(color) != 6:
                return False
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
            return luminance < 0.5

        def is_light(color: str) -> bool:
            """Check if color is light (good for background)."""
            color = color.upper().lstrip('#')
            if len(color) != 6:
                return False
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
            return luminance > 0.7

        # Combine all colors and sort by frequency
        all_colors = fill_colors + text_colors + bg_colors
        sorted_colors = [c for c, _ in all_colors.most_common()]

        # Filter out neutrals for accent colors
        accent_colors = [c for c in sorted_colors if not is_neutral(c)]

        # Find best colors for each role
        primary = accent_colors[0] if len(accent_colors) > 0 else None
        secondary = accent_colors[1] if len(accent_colors) > 1 else None
        accent = accent_colors[2] if len(accent_colors) > 2 else None

        # Find text color (prefer dark colors from text_colors)
        text_sorted = [c for c, _ in text_colors.most_common()]
        text_color = None
        for c in text_sorted:
            if is_dark(c):
                text_color = c
                break
        if not text_color and text_sorted:
            text_color = text_sorted[0]

        # Find background color (prefer light colors from bg_colors)
        bg_sorted = [c for c, _ in bg_colors.most_common()]
        background = None
        for c in bg_sorted:
            if is_light(c):
                background = c
                break
        # Also check fill colors for background
        if not background:
            for c in sorted_colors:
                if is_light(c):
                    background = c
                    break

        return ExtractedColorPalette(
            primary=primary,
            secondary=secondary,
            accent=accent,
            background=background,
            text=text_color,
            all_colors=sorted_colors[:20],  # Top 20 colors
        )

    def _extract_background_images(self) -> list[str]:
        """Extract any embedded background images from the template."""
        extracted = []

        try:
            # Open PPTX as zip to access media files
            with ZipFile(io.BytesIO(self.file_content)) as zf:
                for name in zf.namelist():
                    if name.startswith("ppt/media/") and name.lower().endswith(
                        (".png", ".jpg", ".jpeg", ".gif")
                    ):
                        # Extract to output directory
                        img_data = zf.read(name)
                        img_name = Path(name).name
                        output_path = self.output_dir / f"template_{img_name}"
                        output_path.write_bytes(img_data)
                        extracted.append(str(output_path))
        except Exception:
            pass

        return extracted


async def extract_template_styles(
    file_content: bytes, filename: str, output_dir: str
) -> TemplateProfile:
    """Convenience function to extract template styles from PPTX."""
    extractor = PptxTemplateExtractor(file_content, filename, output_dir)
    return extractor.extract()
