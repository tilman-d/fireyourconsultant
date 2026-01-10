"""PowerPoint generator using python-pptx to create professional PPTX files."""

from pathlib import Path
from typing import Optional
import io

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from PIL import Image

from ..models import (
    BrandProfile,
    Presentation,
    SlideContent,
    SlideLayout,
    ImageCategory,
    ScrapedImage,
    Stat,
)


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = ''.join([c*2 for c in hex_color])
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def lighten_color(hex_color: str, factor: float = 0.3) -> str:
    """Lighten a hex color by a factor."""
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)

    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)

    return f"#{r:02x}{g:02x}{b:02x}"


def darken_color(hex_color: str, factor: float = 0.2) -> str:
    """Darken a hex color by a factor."""
    hex_color = hex_color.lstrip("#")
    r = int(int(hex_color[0:2], 16) * (1 - factor))
    g = int(int(hex_color[2:4], 16) * (1 - factor))
    b = int(int(hex_color[4:6], 16) * (1 - factor))
    return f"#{r:02x}{g:02x}{b:02x}"


class PptxGenerator:
    """Generates professional PowerPoint presentations."""

    # Slide dimensions (16:9 widescreen)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Professional margins
    MARGIN_LEFT = Inches(0.75)
    MARGIN_RIGHT = Inches(0.75)
    MARGIN_TOP = Inches(0.6)
    MARGIN_BOTTOM = Inches(0.5)

    CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    CONTENT_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP - MARGIN_BOTTOM

    def __init__(self, brand: BrandProfile):
        self.brand = brand
        self.prs = PptxPresentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT

        # Image lookup by category
        self.images_by_category: dict[ImageCategory, list[ScrapedImage]] = {}
        self._build_image_lookup()

        # Track which images have been used
        self.used_images: set[str] = set()

        # Slide counter for alternating styles
        self.slide_count = 0

    def _build_image_lookup(self):
        """Build lookup dictionary of images by category."""
        for img in self.brand.images:
            if img.category not in self.images_by_category:
                self.images_by_category[img.category] = []
            self.images_by_category[img.category].append(img)

    def generate(self, presentation: Presentation, output_path: str) -> str:
        """Generate a PowerPoint file from the presentation content."""
        for slide_content in presentation.slides:
            self._add_slide(slide_content)
            self.slide_count += 1

        # Save the presentation
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))

        return str(output_path)

    def _add_slide(self, content: SlideContent) -> None:
        """Add a slide based on its layout type."""
        layout_handlers = {
            SlideLayout.TITLE: self._add_title_slide,
            SlideLayout.BULLETS: self._add_bullet_slide,
            SlideLayout.TWO_COLUMN: self._add_two_column_slide,
            SlideLayout.IMAGE_LEFT: self._add_image_left_slide,
            SlideLayout.IMAGE_RIGHT: self._add_image_right_slide,
            SlideLayout.SECTION_DIVIDER: self._add_section_divider,
            SlideLayout.QUOTE: self._add_quote_slide,
            SlideLayout.STATS: self._add_stats_slide,
            SlideLayout.THANK_YOU: self._add_thank_you_slide,
        }

        handler = layout_handlers.get(content.layout, self._add_bullet_slide)
        handler(content)

    def _create_blank_slide(self):
        """Create a blank slide with white background."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        # Set white background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        return slide

    def _add_gradient_shape(self, slide, left, top, width, height, color1: str, color2: str, angle: int = 90):
        """Add a shape with gradient fill."""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.line.fill.background()

        fill = shape.fill
        fill.gradient()
        fill.gradient_angle = angle
        fill.gradient_stops[0].color.rgb = hex_to_rgb(color1)
        fill.gradient_stops[1].color.rgb = hex_to_rgb(color2)

        return shape

    def _add_accent_bar(self, slide, top: float = Inches(0)):
        """Add a thin accent bar at the top of the slide."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, top,
            self.SLIDE_WIDTH, Inches(0.08)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent)
        shape.line.fill.background()
        return shape

    def _add_title_slide(self, content: SlideContent) -> None:
        """Create a stunning title slide."""
        slide = self._create_blank_slide()

        # Full-width gradient background
        primary = self.brand.colors.primary
        secondary = darken_color(primary, 0.3)

        self._add_gradient_shape(
            slide, 0, 0,
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
            primary, secondary, 135
        )

        # Decorative accent shape in corner
        accent_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            self.SLIDE_WIDTH - Inches(4), Inches(-2),
            Inches(6), Inches(6)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent)
        accent_shape.fill.fore_color.brightness = 0.2
        accent_shape.line.fill.background()

        # Add logo if available
        if self.brand.logo_path and Path(self.brand.logo_path).exists():
            self._add_logo(slide, self.MARGIN_LEFT, Inches(0.5), max_height=Inches(0.9))

        # Main title
        title_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT,
            Inches(2.4),
            self.CONTENT_WIDTH,
            Inches(2),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = self.brand.fonts.heading
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.1

        # Subtitle with accent color underline effect
        if content.subtitle:
            subtitle_box = slide.shapes.add_textbox(
                self.MARGIN_LEFT,
                Inches(4.6),
                Inches(8),
                Inches(1),
            )
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content.subtitle
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.color.brightness = 0.2
            p.font.name = self.brand.fonts.body
            p.alignment = PP_ALIGN.LEFT

        # Bottom accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.MARGIN_LEFT, Inches(6.8),
            Inches(2), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent)
        line.line.fill.background()

        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _add_bullet_slide(self, content: SlideContent) -> None:
        """Create a visually engaging bullet point slide with alternating styles."""
        slide = self._create_blank_slide()

        # Alternate between two visual styles based on slide count
        use_accent_style = (self.slide_count % 2 == 0)

        if use_accent_style:
            # Style A: Subtle gradient background with accent shapes
            self._add_gradient_shape(
                slide, 0, 0,
                self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
                "#ffffff", lighten_color(self.brand.colors.primary, 0.95), 180
            )
            # Decorative corner shape
            corner = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                self.SLIDE_WIDTH - Inches(3), Inches(-1.5),
                Inches(5), Inches(5)
            )
            corner.fill.solid()
            corner.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent)
            corner.fill.fore_color.brightness = 0.7
            corner.line.fill.background()
        else:
            # Style B: Clean white with bold left accent
            # Left accent bar (thicker)
            strip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0,
                Inches(0.25), self.SLIDE_HEIGHT
            )
            strip.fill.solid()
            strip.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.primary)
            strip.line.fill.background()

        # Top accent bar
        self._add_accent_bar(slide)

        # Title with underline accent
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5),
            Inches(11), Inches(1),
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.brand.colors.text)
        p.font.name = self.brand.fonts.heading

        # Title underline accent
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1.25),
            Inches(1.5), Inches(0.05)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent)
        underline.line.fill.background()

        # Subtitle line
        if content.subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.45),
                Inches(10), Inches(0.5),
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = content.subtitle
            p.font.size = Pt(16)
            p.font.color.rgb = hex_to_rgb(self.brand.colors.text_light)
            p.font.name = self.brand.fonts.body

        # Bullets with enhanced styling
        bullet_start_y = Inches(2.2) if content.subtitle else Inches(1.8)
        num_bullets = len(content.bullets[:5])  # Max 5 bullets for cleaner look

        for i, bullet in enumerate(content.bullets[:5]):
            y_pos = bullet_start_y + Inches(i * 0.95)

            # Numbered bullet with background
            num_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.9), y_pos + Inches(0.05),
                Inches(0.4), Inches(0.4)
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent if i == 0 else lighten_color(self.brand.colors.primary, 0.6))
            num_circle.line.fill.background()

            # Number text
            num_box = slide.shapes.add_textbox(
                Inches(0.9), y_pos + Inches(0.07),
                Inches(0.4), Inches(0.4)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Bullet text
            text_box = slide.shapes.add_textbox(
                Inches(1.5), y_pos,
                Inches(10.5), Inches(0.9),
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = hex_to_rgb(self.brand.colors.text)
            p.font.name = self.brand.fonts.body
            p.line_spacing = 1.3

        # Add logo in corner
        if self.brand.logo_path and Path(self.brand.logo_path).exists():
            self._add_logo(slide, self.SLIDE_WIDTH - Inches(1.5), Inches(6.5), max_height=Inches(0.6))

        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _add_two_column_slide(self, content: SlideContent) -> None:
        """Create a professional two-column comparison slide."""
        slide = self._create_blank_slide()

        # Top accent bar
        self._add_accent_bar(slide)

        # Title
        title_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT, Inches(0.5),
            self.CONTENT_WIDTH, Inches(1),
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.brand.colors.text)
        p.font.name = self.brand.fonts.heading

        col_width = Inches(5.5)
        col_height = Inches(5)
        col_top = Inches(1.8)

        # Left column card
        left_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.6), col_top,
            col_width, col_height
        )
        left_card.fill.solid()
        left_card.fill.fore_color.rgb = hex_to_rgb(lighten_color(self.brand.colors.primary, 0.9))
        left_card.line.color.rgb = hex_to_rgb(lighten_color(self.brand.colors.primary, 0.7))
        left_card.line.width = Pt(1)

        # Left column header
        left_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.6), col_top,
            col_width, Inches(0.7)
        )
        left_header.fill.solid()
        left_header.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.primary)
        left_header.line.fill.background()

        # Right column card
        right_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.9), col_top,
            col_width, col_height
        )
        right_card.fill.solid()
        right_card.fill.fore_color.rgb = hex_to_rgb(lighten_color(self.brand.colors.accent, 0.85))
        right_card.line.color.rgb = hex_to_rgb(lighten_color(self.brand.colors.accent, 0.6))
        right_card.line.width = Pt(1)

        # Right column header
        right_header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.9), col_top,
            col_width, Inches(0.7)
        )
        right_header.fill.solid()
        right_header.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent)
        right_header.line.fill.background()

        # Left content
        if content.left_content:
            # Extract header if present
            left_text = content.left_content
            left_header_text = "Option A"
            if left_text.startswith("**") and "**" in left_text[2:]:
                end_idx = left_text.index("**", 2)
                left_header_text = left_text[2:end_idx]
                left_text = left_text[end_idx+2:].strip()

            # Header text
            lh_box = slide.shapes.add_textbox(Inches(0.8), col_top + Inches(0.15), col_width - Inches(0.4), Inches(0.5))
            tf = lh_box.text_frame
            p = tf.paragraphs[0]
            p.text = left_header_text
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = self.brand.fonts.heading
            p.alignment = PP_ALIGN.CENTER

            # Content
            left_box = slide.shapes.add_textbox(
                Inches(0.9), col_top + Inches(0.9),
                col_width - Inches(0.6), col_height - Inches(1.1)
            )
            tf = left_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = left_text.replace("• ", "\n• ").strip()
            p.font.size = Pt(16)
            p.font.color.rgb = hex_to_rgb(self.brand.colors.text)
            p.font.name = self.brand.fonts.body
            p.line_spacing = 1.5

        # Right content
        if content.right_content:
            right_text = content.right_content
            right_header_text = "Option B"
            if right_text.startswith("**") and "**" in right_text[2:]:
                end_idx = right_text.index("**", 2)
                right_header_text = right_text[2:end_idx]
                right_text = right_text[end_idx+2:].strip()

            # Header text
            rh_box = slide.shapes.add_textbox(Inches(7.1), col_top + Inches(0.15), col_width - Inches(0.4), Inches(0.5))
            tf = rh_box.text_frame
            p = tf.paragraphs[0]
            p.text = right_header_text
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = self.brand.fonts.heading
            p.alignment = PP_ALIGN.CENTER

            # Content
            right_box = slide.shapes.add_textbox(
                Inches(7.2), col_top + Inches(0.9),
                col_width - Inches(0.6), col_height - Inches(1.1)
            )
            tf = right_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = right_text.replace("• ", "\n• ").strip()
            p.font.size = Pt(16)
            p.font.color.rgb = hex_to_rgb(self.brand.colors.text)
            p.font.name = self.brand.fonts.body
            p.line_spacing = 1.5

        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _add_image_left_slide(self, content: SlideContent) -> None:
        """Create a slide with image on the left, content on right."""
        slide = self._create_blank_slide()

        # Image area with brand color overlay
        img_width = Inches(5.5)

        # Background for image area
        img_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            img_width, self.SLIDE_HEIGHT
        )
        img_bg.fill.solid()
        img_bg.fill.fore_color.rgb = hex_to_rgb(lighten_color(self.brand.colors.primary, 0.85))
        img_bg.line.fill.background()

        # Add image if available
        image = self._get_image_for_category(content.image_category)
        if image and image.local_path and Path(image.local_path).exists():
            self._add_image_to_slide(
                slide, image.local_path,
                Inches(0.3), Inches(0.5),
                img_width - Inches(0.6), Inches(6.5)
            )

        # Accent bar on right side of image
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            img_width - Inches(0.08), 0,
            Inches(0.08), self.SLIDE_HEIGHT
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent)
        accent.line.fill.background()

        # Content area
        content_left = Inches(6)
        content_width = Inches(6.5)

        # Title
        title_box = slide.shapes.add_textbox(
            content_left, Inches(1),
            content_width, Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.brand.colors.text)
        p.font.name = self.brand.fonts.heading

        # Content
        content_text = content.body_text or "\n".join(f"• {b}" for b in content.bullets)
        if content_text:
            body_box = slide.shapes.add_textbox(
                content_left, Inches(2.4),
                content_width, Inches(4.5)
            )
            tf = body_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content_text
            p.font.size = Pt(18)
            p.font.color.rgb = hex_to_rgb(self.brand.colors.text)
            p.font.name = self.brand.fonts.body
            p.line_spacing = 1.6

        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _add_image_right_slide(self, content: SlideContent) -> None:
        """Create a slide with content on left, image on right."""
        slide = self._create_blank_slide()

        # Top accent bar
        self._add_accent_bar(slide)

        img_width = Inches(5.5)
        img_left = self.SLIDE_WIDTH - img_width

        # Background for image area
        img_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            img_left, 0,
            img_width, self.SLIDE_HEIGHT
        )
        img_bg.fill.solid()
        img_bg.fill.fore_color.rgb = hex_to_rgb(lighten_color(self.brand.colors.secondary, 0.85))
        img_bg.line.fill.background()

        # Add image
        image = self._get_image_for_category(content.image_category)
        if image and image.local_path and Path(image.local_path).exists():
            self._add_image_to_slide(
                slide, image.local_path,
                img_left + Inches(0.3), Inches(0.5),
                img_width - Inches(0.6), Inches(6.5)
            )

        # Content area
        content_width = Inches(6.5)

        # Title
        title_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT, Inches(1),
            content_width, Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.brand.colors.text)
        p.font.name = self.brand.fonts.heading

        # Content
        content_text = content.body_text or "\n".join(f"• {b}" for b in content.bullets)
        if content_text:
            body_box = slide.shapes.add_textbox(
                self.MARGIN_LEFT, Inches(2.4),
                content_width, Inches(4.5)
            )
            tf = body_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content_text
            p.font.size = Pt(18)
            p.font.color.rgb = hex_to_rgb(self.brand.colors.text)
            p.font.name = self.brand.fonts.body
            p.line_spacing = 1.6

        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _add_section_divider(self, content: SlideContent) -> None:
        """Create an impactful section divider slide."""
        slide = self._create_blank_slide()

        # Full gradient background
        primary = self.brand.colors.secondary
        darker = darken_color(primary, 0.4)

        self._add_gradient_shape(
            slide, 0, 0,
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
            primary, darker, 120
        )

        # Large decorative number/shape
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-1), Inches(-1),
            Inches(5), Inches(5)
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent)
        num_shape.fill.fore_color.brightness = 0.3
        num_shape.line.fill.background()

        # Section title - large and centered
        title_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT, Inches(2.8),
            self.CONTENT_WIDTH, Inches(2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = self.brand.fonts.heading
        p.alignment = PP_ALIGN.CENTER

        # Subtle subtitle if present
        if content.subtitle:
            sub_box = slide.shapes.add_textbox(
                self.MARGIN_LEFT, Inches(5),
                self.CONTENT_WIDTH, Inches(0.8)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = content.subtitle
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.color.brightness = 0.3
            p.font.name = self.brand.fonts.body
            p.alignment = PP_ALIGN.CENTER

        # Bottom accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.5), Inches(6.5),
            Inches(2.333), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent)
        line.line.fill.background()

        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _add_quote_slide(self, content: SlideContent) -> None:
        """Create an elegant quote/testimonial slide."""
        slide = self._create_blank_slide()

        # Light background with subtle gradient
        bg_shape = self._add_gradient_shape(
            slide, 0, 0,
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
            lighten_color(self.brand.colors.primary, 0.95),
            "#ffffff", 180
        )

        # Large quotation mark
        quote_mark = slide.shapes.add_textbox(
            Inches(1), Inches(1),
            Inches(2), Inches(2)
        )
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = "\u201C"
        p.font.size = Pt(200)
        p.font.color.rgb = hex_to_rgb(self.brand.colors.accent)
        p.font.color.brightness = 0.5
        p.font.name = "Georgia"

        # Quote text
        quote_box = slide.shapes.add_textbox(
            Inches(2), Inches(2.5),
            Inches(9), Inches(3)
        )
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.quote or content.body_text or content.title
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.color.rgb = hex_to_rgb(self.brand.colors.text)
        p.font.name = self.brand.fonts.body
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.5

        # Author attribution
        if content.quote_author:
            author_box = slide.shapes.add_textbox(
                Inches(2), Inches(5.8),
                Inches(9), Inches(0.6)
            )
            tf = author_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {content.quote_author}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(self.brand.colors.accent)
            p.font.name = self.brand.fonts.body

        # Bottom accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, self.SLIDE_HEIGHT - Inches(0.15),
            self.SLIDE_WIDTH, Inches(0.15)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.primary)
        bar.line.fill.background()

        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _add_stats_slide(self, content: SlideContent) -> None:
        """Create an impressive stats/metrics slide with large numbers."""
        slide = self._create_blank_slide()

        # Subtle gradient background
        self._add_gradient_shape(
            slide, 0, 0,
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
            "#ffffff", lighten_color(self.brand.colors.primary, 0.97), 135
        )

        # Top accent bar
        self._add_accent_bar(slide)

        # Title
        title_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT, Inches(0.5),
            self.CONTENT_WIDTH, Inches(1),
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.brand.colors.text)
        p.font.name = self.brand.fonts.heading

        # Parse stats from bullets if no explicit stats provided
        stats_data = []
        if content.stats:
            stats_data = content.stats
        elif content.bullets:
            # Try to parse bullets as stats (format: "Value - Description")
            for bullet in content.bullets[:4]:
                if " - " in bullet:
                    parts = bullet.split(" - ", 1)
                    stats_data.append(Stat(value=parts[0].strip(), label=parts[1].strip()))
                elif bullet[0].isdigit() or bullet[0] in "$%":
                    # Try to extract number from start
                    import re
                    match = re.match(r'^([\d$%,.]+\w*)\s*(.*)$', bullet)
                    if match:
                        stats_data.append(Stat(value=match.group(1), label=match.group(2)))

        num_stats = min(len(stats_data), 4)
        if num_stats == 0:
            return self._add_bullet_slide(content)  # Fallback

        # Calculate positions
        card_width = Inches(2.7)
        card_height = Inches(3.8)
        total_width = num_stats * card_width + (num_stats - 1) * Inches(0.4)
        start_x = (self.SLIDE_WIDTH - total_width) / 2
        card_top = Inches(2)

        for i, stat in enumerate(stats_data[:4]):
            x_pos = start_x + i * (card_width + Inches(0.4))

            # Card background with rounded corners
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, card_top,
                card_width, card_height
            )
            card.fill.solid()
            # Alternate card colors
            if i % 2 == 0:
                card.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.primary)
            else:
                card.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent)
            card.line.fill.background()

            # Large stat value
            value_box = slide.shapes.add_textbox(
                x_pos, card_top + Inches(0.8),
                card_width, Inches(1.5)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = stat.value if isinstance(stat, Stat) else str(stat.get('value', ''))
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = self.brand.fonts.heading
            p.alignment = PP_ALIGN.CENTER

            # Stat label
            label_box = slide.shapes.add_textbox(
                x_pos + Inches(0.2), card_top + Inches(2.4),
                card_width - Inches(0.4), Inches(1.2)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = stat.label if isinstance(stat, Stat) else str(stat.get('label', ''))
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.color.brightness = 0.2
            p.font.name = self.brand.fonts.body
            p.alignment = PP_ALIGN.CENTER

        # Bottom decorative line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.5), Inches(6.5),
            Inches(2.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.brand.colors.accent)
        line.line.fill.background()

        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _add_thank_you_slide(self, content: SlideContent) -> None:
        """Create a compelling closing slide."""
        slide = self._create_blank_slide()

        # Full gradient background
        primary = self.brand.colors.primary
        accent = self.brand.colors.accent

        self._add_gradient_shape(
            slide, 0, 0,
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
            primary, darken_color(primary, 0.3), 135
        )

        # Decorative circles
        for i, (x, y, size, opacity) in enumerate([
            (10, -1, 4, 0.1),
            (11, 5, 3, 0.08),
            (-1, 5, 3.5, 0.12),
        ]):
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y),
                Inches(size), Inches(size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = hex_to_rgb(accent)
            circle.fill.fore_color.brightness = 0.5 - opacity
            circle.line.fill.background()

        # Logo centered
        if self.brand.logo_path and Path(self.brand.logo_path).exists():
            self._add_logo(slide, Inches(5.5), Inches(1.2), max_height=Inches(1.2))

        # Thank you text
        title_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT, Inches(3),
            self.CONTENT_WIDTH, Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.title or "Thank You"
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = self.brand.fonts.heading
        p.alignment = PP_ALIGN.CENTER

        # Contact info / CTA
        if content.body_text or content.subtitle:
            contact_box = slide.shapes.add_textbox(
                self.MARGIN_LEFT, Inches(4.8),
                self.CONTENT_WIDTH, Inches(1.2)
            )
            tf = contact_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content.body_text or content.subtitle or "Let's connect"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.color.brightness = 0.2
            p.font.name = self.brand.fonts.body
            p.alignment = PP_ALIGN.CENTER

        # Bottom accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5), Inches(6.2),
            Inches(3.333), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(accent)
        line.line.fill.background()

        if content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = content.speaker_notes

    def _add_logo(self, slide, left: float, top: float, max_height: float = Inches(1)) -> None:
        """Add logo maintaining aspect ratio."""
        logo_path = Path(self.brand.logo_path)
        if not logo_path.exists():
            return

        # Skip SVGs - can't be embedded easily
        if logo_path.suffix.lower() == '.svg':
            return

        try:
            with Image.open(logo_path) as img:
                width, height = img.size
                aspect = width / height

                new_height = max_height
                new_width = max_height * aspect

                # Cap width
                if new_width > Inches(3):
                    new_width = Inches(3)
                    new_height = new_width / aspect

                slide.shapes.add_picture(
                    str(logo_path), left, top,
                    width=new_width, height=new_height
                )
        except Exception as e:
            print(f"Error adding logo: {e}")

    def _add_image_to_slide(self, slide, image_path: str, left: float, top: float,
                           max_width: float, max_height: float) -> None:
        """Add image with smart sizing."""
        path = Path(image_path)
        if not path.exists() or path.suffix.lower() == '.svg':
            return

        try:
            with Image.open(path) as img:
                width, height = img.size
                aspect = width / height

                if width / max_width > height / max_height:
                    new_width = max_width
                    new_height = max_width / aspect
                else:
                    new_height = max_height
                    new_width = max_height * aspect

                # Center the image
                x_offset = (max_width - new_width) / 2
                y_offset = (max_height - new_height) / 2

                slide.shapes.add_picture(
                    str(path),
                    left + x_offset, top + y_offset,
                    width=new_width, height=new_height
                )
        except Exception as e:
            print(f"Error adding image: {e}")

    def _get_image_for_category(self, category: Optional[ImageCategory]) -> Optional[ScrapedImage]:
        """Get an unused image for the category, prioritizing user uploads."""
        # Always try user-uploaded images first (highest priority)
        if ImageCategory.USER_UPLOAD in self.images_by_category:
            for img in self.images_by_category[ImageCategory.USER_UPLOAD]:
                if img.local_path and img.local_path not in self.used_images:
                    if not img.local_path.endswith('.svg'):
                        self.used_images.add(img.local_path)
                        return img

        if not category:
            # Try to get any good image
            for cat in [ImageCategory.HERO, ImageCategory.PRODUCT, ImageCategory.TEAM, ImageCategory.OFFICE]:
                if cat in self.images_by_category:
                    for img in self.images_by_category[cat]:
                        if img.local_path and img.local_path not in self.used_images:
                            if not img.local_path.endswith('.svg'):
                                self.used_images.add(img.local_path)
                                return img
            return None

        images = self.images_by_category.get(category, [])

        for img in images:
            if img.local_path and img.local_path not in self.used_images:
                if not img.local_path.endswith('.svg'):
                    self.used_images.add(img.local_path)
                    return img

        # Reuse user uploads before other images
        if ImageCategory.USER_UPLOAD in self.images_by_category:
            for img in self.images_by_category[ImageCategory.USER_UPLOAD]:
                if img.local_path and not img.local_path.endswith('.svg'):
                    return img

        # Reuse other images if needed
        for img in images:
            if img.local_path and not img.local_path.endswith('.svg'):
                return img

        return None


def generate_pptx(presentation: Presentation, brand: BrandProfile, output_path: str) -> str:
    """Generate a PowerPoint file."""
    generator = PptxGenerator(brand)
    return generator.generate(presentation, output_path)
