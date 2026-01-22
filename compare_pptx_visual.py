#!/usr/bin/env python3
"""Compare PPTX files by converting to images and visual comparison."""

import subprocess
from pathlib import Path


def convert_pptx_to_pdf(pptx_path: str, output_dir: str):
    """Convert PPTX to PDF using LibreOffice."""
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # Use LibreOffice to convert
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(pptx_path)
    ]

    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"Error converting {pptx_path}: {result.stderr}")
        return None

    pdf_path = output_dir / f"{pptx_path.stem}.pdf"
    if pdf_path.exists():
        print(f"Created PDF: {pdf_path}")
        return str(pdf_path)
    return None


def analyze_with_pptx(template_path: str, generated_path: str):
    """Analyze both PPTX files in detail."""
    from pptx import Presentation as PptxPresentation

    template = PptxPresentation(template_path)
    generated = PptxPresentation(generated_path)

    print("="*80)
    print("DETAILED COMPARISON")
    print("="*80)

    print(f"\n{'Aspect':<25} {'Template':<30} {'Generated':<30}")
    print("-"*85)
    print(f"{'Slides':<25} {len(template.slides):<30} {len(generated.slides):<30}")

    # Compare layouts used
    template_layouts = [s.slide_layout.name for s in template.slides if s.slide_layout]
    generated_layouts = [s.slide_layout.name for s in generated.slides if s.slide_layout]

    print(f"\n{'Layout Distribution:'}")
    from collections import Counter
    t_counts = Counter(template_layouts)
    g_counts = Counter(generated_layouts)

    all_layouts = set(t_counts.keys()) | set(g_counts.keys())
    for layout in sorted(all_layouts):
        t = t_counts.get(layout, 0)
        g = g_counts.get(layout, 0)
        match = "✓" if (t > 0 and g > 0) else ""
        print(f"  {layout:<30} {t:<10} {g:<10} {match}")

    # Compare first few slides in detail
    print("\n" + "="*80)
    print("SLIDE-BY-SLIDE COMPARISON (first 5 generated slides)")
    print("="*80)

    for i, slide in enumerate(list(generated.slides)[:5]):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        print(f"\n--- Generated Slide {i+1}: {layout_name} ---")

        # Get text content
        for shape in slide.shapes:
            try:
                ph_format = shape.placeholder_format
                idx = ph_format.idx
                ph_type = str(ph_format.type)
            except (ValueError, AttributeError):
                idx = None
                ph_type = None

            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text_frame.text.strip()[:80]
                if text:
                    left = shape.left.inches if shape.left else 0
                    top = shape.top.inches if shape.top else 0
                    print(f"  [ph_idx={idx}] pos=({left:.1f}\", {top:.1f}\") '{text}...'")


if __name__ == "__main__":
    template_path = "/root/fyc/output/58388ab2/uploaded_template_20200417_Scoping Präsentation.pptx"
    generated_path = "/root/fyc/output/test_generation/test_presentation.pptx"

    analyze_with_pptx(template_path, generated_path)
