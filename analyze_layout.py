#!/usr/bin/env python3
"""Detailed layout analysis to understand template structure."""

import io
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt

def analyze_layout_details(filepath: str):
    """Analyze layout placeholders in detail."""
    prs = PptxPresentation(filepath)

    print("="*80)
    print("DETAILED LAYOUT ANALYSIS")
    print("="*80)

    for i, layout in enumerate(prs.slide_layouts):
        print(f"\n{'='*60}")
        print(f"LAYOUT {i}: {layout.name}")
        print(f"{'='*60}")

        print("\nPlaceholders:")
        for ph in layout.placeholders:
            idx = ph.placeholder_format.idx
            ph_type = str(ph.placeholder_format.type).replace("PP_PLACEHOLDER.", "")

            left = ph.left.inches if ph.left else 0
            top = ph.top.inches if ph.top else 0
            width = ph.width.inches if ph.width else 0
            height = ph.height.inches if ph.height else 0

            print(f"  idx={idx}, type={ph_type}")
            print(f"    Position: ({left:.2f}\", {top:.2f}\") Size: {width:.2f}\" x {height:.2f}\"")

            # Check for text styling
            if ph.has_text_frame:
                tf = ph.text_frame
                for para in tf.paragraphs:
                    font_name = para.font.name
                    font_size = para.font.size.pt if para.font.size else None
                    font_bold = para.font.bold
                    text = para.text[:30] if para.text else ""
                    print(f"    Font: {font_name}, Size: {font_size}pt, Bold: {font_bold}")
                    if text:
                        print(f"    Sample text: '{text}'")
                    break  # Just first para

        # Check for non-placeholder shapes in the layout
        def is_placeholder(s):
            try:
                _ = s.placeholder_format
                return True
            except ValueError:
                return False

        non_ph_shapes = [s for s in layout.shapes if not is_placeholder(s)]
        if non_ph_shapes:
            print(f"\n  Non-placeholder shapes: {len(non_ph_shapes)}")
            for shape in non_ph_shapes[:5]:
                print(f"    - {type(shape).__name__}")
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    text = shape.text_frame.text[:50] if shape.text_frame.text else ""
                    if text:
                        print(f"      Text: '{text}'")

    # Also analyze a few sample slides to see how content appears
    print("\n" + "="*80)
    print("SAMPLE SLIDE CONTENT ANALYSIS")
    print("="*80)

    for i, slide in enumerate(list(prs.slides)[:5]):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        print(f"\n--- Slide {i+1} (Layout: {layout_name}) ---")

        for shape in slide.shapes:
            shape_type = type(shape).__name__

            # Check if it's a placeholder
            try:
                ph_format = shape.placeholder_format
                is_placeholder = True
                ph_idx = ph_format.idx
            except (ValueError, AttributeError):
                is_placeholder = False
                ph_idx = None

            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text_frame.text.strip()[:60] if shape.text_frame.text else ""
                if text:
                    # Get styling
                    try:
                        p = shape.text_frame.paragraphs[0]
                        font = p.font.name or "inherit"
                        size = p.font.size.pt if p.font.size else "inherit"
                        print(f"  [{shape_type}, ph_idx={ph_idx}] '{text}' (font={font}, size={size})")
                    except:
                        print(f"  [{shape_type}, ph_idx={ph_idx}] '{text}'")


if __name__ == "__main__":
    template_path = "/root/fyc/output/58388ab2/uploaded_template_20200417_Scoping Präsentation.pptx"
    analyze_layout_details(template_path)

    print("\n\n" + "="*80)
    print("GENERATED PRESENTATION ANALYSIS")
    print("="*80)

    generated_path = "/root/fyc/output/58388ab2/presentation.pptx"
    prs = PptxPresentation(generated_path)

    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        print(f"\n--- Slide {i+1} (Layout: {layout_name}) ---")

        for shape in slide.shapes:
            shape_type = type(shape).__name__
            is_placeholder = hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None
            ph_idx = shape.placeholder_format.idx if is_placeholder else None

            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text_frame.text.strip()[:60] if shape.text_frame.text else ""
                if text:
                    try:
                        p = shape.text_frame.paragraphs[0]
                        font = p.font.name or "inherit"
                        size = p.font.size.pt if p.font.size else "inherit"
                        print(f"  [{shape_type}, ph_idx={ph_idx}] '{text}' (font={font}, size={size})")
                    except:
                        print(f"  [{shape_type}, ph_idx={ph_idx}] '{text}'")
