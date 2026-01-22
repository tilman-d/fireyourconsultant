#!/usr/bin/env python3
"""Test presentation generation with template to verify styling."""

import asyncio
import sys
from pathlib import Path

# Add src to path
sys.path.insert(0, str(Path(__file__).parent / "src"))

from fyc.pptx_gen.generator import PptxGenerator
from fyc.template.extractor import PptxTemplateExtractor
from fyc.models import (
    BrandProfile, BrandColors, BrandFonts, BrandVoice,
    Presentation, SlideContent, SlideLayout, ImageCategory, ScrapedImage
)


def test_generation():
    """Test presentation generation with template."""
    # Load template
    template_path = "/root/fyc/output/58388ab2/uploaded_template_20200417_Scoping Präsentation.pptx"
    output_dir = "/root/fyc/output/test_generation"

    # Create output directory
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    # Read template bytes
    with open(template_path, "rb") as f:
        template_bytes = f.read()

    # Extract template profile
    extractor = PptxTemplateExtractor(template_bytes, "template.pptx", output_dir)
    template_profile = extractor.extract()

    print("Template Profile:")
    print(f"  Theme Colors: accent1={template_profile.theme_colors.accent1}")
    print(f"  Theme Fonts: major={template_profile.theme_fonts.major_latin}, minor={template_profile.theme_fonts.minor_latin}")
    print(f"  Layouts: {[l.name for l in template_profile.layouts]}")

    # Create brand profile
    brand = BrandProfile(
        colors=BrandColors(
            primary="#048EFE",
            secondary="#002F86",
            accent="#FFC000",
            text="#000000",
            text_light="#777776",
            background="#FFFFFF",
        ),
        fonts=BrandFonts(
            heading="Roboto Medium",
            body="Roboto Light",
        ),
        voice=BrandVoice(
            formality=0.8,
            technicality=0.7,
            enthusiasm=0.6,
        ),
        company_name="Test Company",
        tagline="Let us grow together",
        language="de",
        images=[
            # Use extracted template images for content slides
            ScrapedImage(
                url="local",
                local_path=f"{output_dir}/template_image3.jpg",
                category=ImageCategory.PRODUCT,
            ),
            ScrapedImage(
                url="local",
                local_path=f"{output_dir}/template_image4.jpg",
                category=ImageCategory.HERO,
            ),
            ScrapedImage(
                url="local",
                local_path=f"{output_dir}/template_image5.jpg",
                category=ImageCategory.TEAM,
            ),
        ],
        logo_path=None,
        template_profile=template_profile,
    )

    # Create test presentation content
    presentation = Presentation(
        title="Test Presentation",
        slides=[
            SlideContent(
                layout=SlideLayout.TITLE,
                title="AI Strategy Workshop",
                subtitle="Transforming Your Business with Artificial Intelligence",
            ),
            SlideContent(
                layout=SlideLayout.BULLETS,
                title="Key Discussion Points",
                bullets=[
                    "Understanding AI capabilities",
                    "Identifying use cases",
                    "Implementation roadmap",
                    "Resource requirements",
                ],
            ),
            SlideContent(
                layout=SlideLayout.IMAGE_RIGHT,
                title="Market Analysis",
                bullets=[
                    "Growing adoption across industries",
                    "Cost reduction opportunities",
                    "Competitive advantages",
                ],
                image_category=ImageCategory.PRODUCT,
            ),
            SlideContent(
                layout=SlideLayout.STATS,
                title="Expected Results",
                bullets=[
                    "40% - Efficiency improvement",
                    "25% - Cost reduction",
                    "60% - Customer satisfaction increase",
                ],
            ),
            SlideContent(
                layout=SlideLayout.SECTION_DIVIDER,
                title="Implementation",
                subtitle="Phase-by-phase approach",
            ),
            SlideContent(
                layout=SlideLayout.THANK_YOU,
                title="Vielen Dank!",
                body_text="Kontaktieren Sie uns für weitere Informationen",
            ),
        ],
    )

    # Generate presentation
    generator = PptxGenerator(brand)
    output_path = f"{output_dir}/test_presentation.pptx"
    result = generator.generate(presentation, output_path)

    print(f"\nGenerated presentation: {result}")
    print("\nAnalyzing generated presentation...")

    # Analyze the generated presentation
    from pptx import Presentation as PptxPresentation
    prs = PptxPresentation(result)

    print(f"\nGenerated {len(prs.slides)} slides:")
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        print(f"  Slide {i+1}: Layout = '{layout_name}'")

        # Check for text content
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text_frame.text.strip()[:50]
                if text:
                    print(f"    - Text: '{text}...'")

    return result


if __name__ == "__main__":
    test_generation()
